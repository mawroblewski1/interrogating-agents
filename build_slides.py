#!/usr/bin/env python3
"""Build slides 2-7 for ECS 172 presentation, preserving slide 1."""

from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

PPTX = r"D:\Academic\2026\ECS_172\Project\ECS 172 Group Project Presentation - Moderating Extremists via Interrogation.pptx"
FIG1 = r"D:\Academic\2026\ECS_172\Project\results\figures\fig1_magnitude_by_condition.png"
FIG3 = r"D:\Academic\2026\ECS_172\Project\results\figures\fig3_directional_accuracy.png"

W, H = 9144000, 5143500

BG    = RGBColor(0x43, 0x43, 0x43)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xE8, 0xA0, 0x20)
NAVY  = RGBColor(0x1A, 0x37, 0x5E)
STEEL = RGBColor(0x2E, 0x75, 0xB6)
RUST  = RGBColor(0xB8, 0x4A, 0x10)
DARK  = RGBColor(0x2D, 0x2D, 0x2D)
LGREY = RGBColor(0xBB, 0xBB, 0xBB)
TEAL  = RGBColor(0x1A, 0x6B, 0x72)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(1.5)):
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as MAST
    s = slide.shapes.add_shape(MAST.RECTANGLE, Emu(x), Emu(y), Emu(w), Emu(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def add_textbox(slide, lines, x, y, w, h, sz=15, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    """lines: str (split on \\n), list of str, or list of (str, bold, color)."""
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if isinstance(lines, str):
        lines = lines.split('\n')
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            text, lbold, lcolor = line
        else:
            text, lbold, lcolor = line, bold, color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if text:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(sz)
            r.font.bold = lbold
            r.font.italic = italic
            r.font.color.rgb = lcolor
    return tb


def add_connector(slide, x1, y1, x2, y2, color=WHITE, arrowhead=True):
    try:
        from pptx.enum.shapes import MSO_CONNECTOR_TYPE
        ctype = MSO_CONNECTOR_TYPE.STRAIGHT
    except Exception:
        ctype = 1
    c = slide.shapes.add_connector(ctype, Emu(x1), Emu(y1), Emu(x2), Emu(y2))
    c.line.color.rgb = color
    c.line.width = Pt(2)
    if arrowhead:
        try:
            spPr = c._element.find(qn('p:spPr'))
            if spPr is not None:
                ln = spPr.find(qn('a:ln'))
                if ln is None:
                    ln = etree.SubElement(spPr, qn('a:ln'))
                tail = etree.SubElement(ln, qn('a:tailEnd'))
                tail.set('type', 'arrow')
                tail.set('w', 'med')
                tail.set('len', 'med')
        except Exception:
            pass
    return c


def slide_header(slide, title):
    add_rect(slide, 0, 0, W, 690000, fill=DARK)
    add_rect(slide, 0, 680000, W, 50000, fill=AMBER)
    add_textbox(slide, title, 300000, 90000, W - 400000, 560000,
                sz=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def clear_slide(slide):
    """Remove all shapes from a slide without touching the part references."""
    sp_tree = slide.shapes._spTree
    for sp in list(sp_tree)[2:]:  # first 2 children are nvGrpSpPr + grpSpPr
        sp_tree.remove(sp)
    # Reset background to solid fill
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


# ── Slide 3: Architecture ─────────────────────────────────────────────────────
def build_slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    set_bg(slide)
    slide_header(slide, "System Architecture")

    BOX_W, BOX_H, ARR_W, BOX_Y = 2320000, 960000, 380000, 810000
    total = 3 * BOX_W + 2 * ARR_W
    pad = (W - total) // 2

    roles = [
        ("SUSPECT",       ["Holds initial", "pro / con stance"],         STEEL),
        ("INTERROGATOR",  ["2-stage RAG", "pipeline"],                   RUST),
        ("JUDGE",         ["Scores stance", "−2 to +2 / turn",
                           "across 6 turns"],                            TEAL),
    ]
    for i, (name, sub, color) in enumerate(roles):
        bx = pad + i * (BOX_W + ARR_W)
        add_rect(slide, bx, BOX_Y, BOX_W, BOX_H, fill=color, line_color=WHITE)
        add_textbox(slide, name, bx + 60000, BOX_Y + 140000, BOX_W - 120000, 340000,
                    sz=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, sub, bx + 60000, BOX_Y + 500000, BOX_W - 120000, 420000,
                    sz=13, color=LGREY, align=PP_ALIGN.CENTER)

    for i in range(2):
        ax = pad + (i + 1) * BOX_W + i * ARR_W + 30000
        ay = BOX_Y + BOX_H // 2
        add_connector(slide, ax, ay, ax + ARR_W - 60000, ay, color=WHITE)

    RAG_Y = BOX_Y + BOX_H + 260000
    RAG_W, RAG_H, RAG_GAP = 2050000, 820000, 220000
    int_cx = pad + BOX_W + ARR_W + BOX_W // 2
    rag_total = 2 * RAG_W + RAG_GAP
    rag_x1 = int_cx - rag_total // 2
    rag_x2 = rag_x1 + RAG_W + RAG_GAP

    rag_panels = [
        ("Technique RAG",
         ["18 cards  (Reid · Second Sight · Debate)",
          "Phase-aware retrieval",
          "ChromaDB + MiniLM-L6-v2"],
         NAVY),
        ("Topic / Argument RAG",
         ["3 ballot-measure topics",
          "Housing · Arts · Transit",
          "Side-specific arguments"],
         RGBColor(0x2A, 0x47, 0x18)),
    ]
    for j, (title, sub, color) in enumerate(rag_panels):
        rx = rag_x1 if j == 0 else rag_x2
        add_rect(slide, rx, RAG_Y, RAG_W, RAG_H, fill=color, line_color=AMBER)
        add_textbox(slide, title, rx + 60000, RAG_Y + 60000, RAG_W - 120000, 280000,
                    sz=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
        add_textbox(slide, sub, rx + 60000, RAG_Y + 350000, RAG_W - 120000, 430000,
                    sz=12, color=LGREY, align=PP_ALIGN.CENTER)
        rcx = rx + RAG_W // 2
        add_connector(slide, rcx, RAG_Y, rcx, BOX_Y + BOX_H + 20000, color=AMBER)

    quad_y = RAG_Y + RAG_H + 180000
    add_textbox(
        slide,
        "Quad structure: 4 matched legs per topic  ·  "
        "2 control (free-form)  +  2 treatment (RAG-guided)  ·  6 turns each",
        300000, quad_y, W - 600000, 350000,
        sz=13, color=LGREY, align=PP_ALIGN.CENTER, italic=True
    )


# ── Slide 4: Experimental Setup ───────────────────────────────────────────────
def build_slide4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    set_bg(slide)
    slide_header(slide, "Experimental Setup")

    TW, TH, GAP, TY = 2600000, 1380000, 172000, 810000
    total = 3 * TW + 2 * GAP
    pad = (W - total) // 2

    topics = [
        ("Housing",  "Rezoning ballot measure",
         ['"Should the city', 'rezoning plan pass?"']),
        ("Arts",     "Arts funding ballot measure",
         ['"Should arts', 'funding increase?"']),
        ("Transit",  "Fare-free transit ballot",
         ['"Should transit', 'fares be eliminated?"']),
    ]
    for i, (name, sub, q) in enumerate(topics):
        tx = pad + i * (TW + GAP)
        add_rect(slide, tx, TY, TW, TH,
                 fill=RGBColor(0x28, 0x28, 0x45), line_color=STEEL, line_w=Pt(1.5))
        add_textbox(slide, name, tx + 60000, TY + 60000, TW - 120000, 300000,
                    sz=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, sub, tx + 60000, TY + 380000, TW - 120000, 250000,
                    sz=13, color=LGREY, align=PP_ALIGN.CENTER, italic=True)
        add_textbox(slide, q, tx + 80000, TY + 670000, TW - 160000, 650000,
                    sz=14, color=AMBER, align=PP_ALIGN.CENTER)

    stats_y = TY + TH + 280000
    stats = [
        ("Model",      ["llama3.1:8b", "(all 3 roles)", "via Ollama"]),
        ("Scale",      ["5 quads × 3 topics", "4 legs × 6 turns",
                        "= 360 scored turns"]),
        ("Techniques", ["18 cards", "Reid (×5) · Second Sight (×7)",
                        "Structured Debate (×6)"]),
        ("Retrieval",  ["all-MiniLM-L6-v2", "ChromaDB index",
                        "Phase-aware"]),
    ]
    sw = (W - 600000) // 4
    for i, (label, vals) in enumerate(stats):
        sx = 300000 + i * sw
        add_rect(slide, sx, stats_y, sw - 80000, 1750000,
                 fill=DARK, line_color=RGBColor(0x55, 0x55, 0x55))
        add_textbox(slide, label, sx + 50000, stats_y + 70000, sw - 180000, 280000,
                    sz=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
        add_textbox(slide, vals, sx + 50000, stats_y + 390000, sw - 180000, 1280000,
                    sz=13, color=WHITE, align=PP_ALIGN.CENTER)


# ── Slide 5: Results — Directional Accuracy ───────────────────────────────────
def build_slide5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    set_bg(slide)
    slide_header(slide, "Results: Directional Accuracy")

    add_rect(slide, 0, 700000, W, 270000, fill=DARK)
    add_textbox(
        slide,
        "Pooled: 77% (treatment) vs 42% (control)   ·   "
        "Transit: 90% vs 10%   ·   "
        "Did the suspect’s stance move toward the interrogator’s target?",
        300000, 715000, W - 600000, 230000,
        sz=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER
    )
    slide.shapes.add_picture(FIG3, Emu(300000), Emu(1010000),
                             Emu(W - 600000), Emu(H - 1110000))


# ── Slide 6: Results — Magnitude ─────────────────────────────────────────────
def build_slide6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    set_bg(slide)
    slide_header(slide, "Results: Stance-Change Magnitude")

    add_rect(slide, 0, 700000, W, 270000, fill=DARK)
    add_textbox(
        slide,
        "How large was the stance shift?  ( |s_final − s_initial| )   ·   "
        "Pooled: treatment 1.83 vs control 1.66   ·   ⚠ Housing confounded by judge bias",
        300000, 715000, W - 600000, 230000,
        sz=14, color=LGREY, align=PP_ALIGN.CENTER
    )
    slide.shapes.add_picture(FIG1, Emu(300000), Emu(1010000),
                             Emu(W - 600000), Emu(H - 1110000))

    # Bias callout overlaid bottom-right
    cx, cy = W - 2850000, H - 980000
    add_rect(slide, cx, cy, 2650000, 820000,
             fill=RGBColor(0x45, 0x18, 0x08), line_color=AMBER)
    add_textbox(slide, [
        "⚠  Housing: Judge Bias",
        "llama3.1:8b scores pro-rezoning statements",
        "negative in batch mode — magnitude metric",
        "is less reliable here than directional accuracy.",
    ], cx + 70000, cy + 70000, 2510000, 700000, sz=12, color=AMBER)


# ── Slide 7: Discussion & Conclusion ─────────────────────────────────────────
def build_slide7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[10])
    set_bg(slide)
    slide_header(slide, "Discussion & Conclusion")

    half = (W - 500000) // 2
    lx, rx = 200000, 200000 + half + 100000
    col_y, col_h = 810000, H - 960000

    add_rect(slide, lx, col_y, half, col_h,
             fill=RGBColor(0x1A, 0x2E, 0x4A), line_color=STEEL)
    add_textbox(slide, "Key Findings",
                lx + 80000, col_y + 70000, half - 160000, 310000,
                sz=19, bold=True, color=STEEL)
    add_textbox(slide, [
        "• RAG-guided interrogation improves",
        "  directional accuracy by +35 pp pooled",
        "  (77% vs 42%)",
        "",
        "• Transit: strongest signal",
        "  (90% treatment vs 10% control)",
        "",
        "• Treatment moves stance in the intended",
        "  direction across all 3 topics",
        "",
        "• RecSys framing opens a new avenue",
        "  for computational moderation",
    ], lx + 80000, col_y + 420000, half - 160000, col_h - 490000, sz=14, color=WHITE)

    add_rect(slide, rx, col_y, half, col_h,
             fill=RGBColor(0x2D, 0x18, 0x0E), line_color=RUST)
    add_textbox(slide, "Limitations & Future Work",
                rx + 80000, col_y + 70000, half - 160000, 310000,
                sz=19, bold=True, color=AMBER)
    add_textbox(slide, [
        "Limitations:",
        "• Single model family (llama3.1:8b) —",
        "  shared bias between suspect & judge",
        "• Small N: 5 quads / topic (indicative)",
        "• LLM-vs-LLM simulation only",
        "",
        "Future Work:",
        "→ Cross-family judge (e.g. qwen2.5:7b)",
        "→ Log technique selections for",
        "   co-occurrence analysis",
        "→ Human subject evaluation",
    ], rx + 80000, col_y + 420000, half - 160000, col_h - 490000, sz=14, color=WHITE)


# ── Main ──────────────────────────────────────────────────────────────────────
def build_slide2_into(prs, slide):
    """Populate slide 2 in place (avoids zip duplicate warning from delete+add)."""
    clear_slide(slide)
    set_bg(slide)
    slide_header(slide, "Motivation & Research Question")

    PW, PH, GAP, PY = 2620000, 3500000, 130000, 810000
    total = 3 * PW + 2 * GAP
    pad = (W - total) // 2

    panels = [
        ("The Problem", NAVY, [
            "Extremist content spreads faster",
            "than humans can moderate.",
            "",
            "Skilled interrogators follow",
            "structured persuasion playbooks",
            "— LLMs currently have none.",
        ]),
        ("The Insight", RGBColor(0x4A, 0x32, 0x0C), [
            "Recommender Systems can",
            "retrieve the right technique",
            "for the right moment.",
            "",
            "Real interrogation corpus",
            "→ RAG → LLM Interrogator",
        ]),
        ("Research Question", RGBColor(0x1A, 0x45, 0x2C), [
            "Can RAG-guided interrogation",
            "reliably shift an LLM's stance",
            "on polarizing topics?",
            "",
            "Metrics: magnitude,",
            "directional accuracy,",
            "consistency",
        ]),
    ]

    for i, (header, fill, body) in enumerate(panels):
        px = pad + i * (PW + GAP)
        add_rect(slide, px, PY, PW, PH, fill=fill, line_color=AMBER)
        add_rect(slide, px, PY, PW, 310000, fill=AMBER)
        add_textbox(slide, header, px + 60000, PY + 55000, PW - 120000, 240000,
                    sz=17, bold=True, color=DARK, align=PP_ALIGN.CENTER)
        add_textbox(slide, body, px + 80000, PY + 380000, PW - 160000, PH - 440000,
                    sz=15, color=WHITE)

    for i in range(2):
        ax = pad + (i + 1) * PW + i * GAP + 10000
        ay = PY + PH // 2
        add_connector(slide, ax, ay, ax + GAP - 20000, ay, color=AMBER)


def main():
    prs = Presentation(PPTX)
    # Rebuild slide 2 in place, then append slides 3-7
    build_slide2_into(prs, prs.slides[1])
    build_slide3(prs)
    build_slide4(prs)
    build_slide4(prs)
    build_slide5(prs)
    build_slide6(prs)
    build_slide7(prs)
    prs.save(PPTX)
    print(f"Done — {len(prs.slides)} slides saved to:\n{PPTX}")


if __name__ == "__main__":
    main()
